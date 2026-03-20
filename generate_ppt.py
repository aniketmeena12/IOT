#!/usr/bin/env python3
"""
Generate a professional client-facing PowerPoint presentation:
  • IoT-based OEE (Overall Equipment Effectiveness) measurement
  • Supply Chain Management, Communication & Automated Goods Tracking
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colours ────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT_BLUE  = RGBColor(0x00, 0x9A, 0xF0)   # bright blue
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)   # green
ACCENT_ORANGE= RGBColor(0xF5, 0x9E, 0x0B)   # orange
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xBB, 0xBB, 0xBB)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
CARD_BG      = RGBColor(0x1A, 0x25, 0x3C)   # card background
HIGHLIGHT    = RGBColor(0x00, 0xD4, 0xAA)   # teal highlight

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ═════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def add_bg(slide, color=DARK_BG):
    """Fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def tb(slide, left, top, width, height):
    """Add a text box and return it."""
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(textbox, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    return p

def add_bullet_list(slide, left, top, width, height, items, size=16, color=WHITE, title=None, title_size=22, title_color=ACCENT_BLUE):
    box = tb(slide, left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = "Calibri"
        first = False
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(8)
    return box

def section_divider(title_text, subtitle_text, accent=ACCENT_BLUE):
    """Full-screen section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    # accent bar
    add_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Pt(4), accent)
    # title
    t = tb(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2))
    set_text(t, title_text, size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # subtitle
    t = tb(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1))
    set_text(t, subtitle_text, size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return slide

def card(slide, left, top, width, height, title, body_lines, icon_text="", accent=ACCENT_BLUE):
    """A styled card with title, icon, and bullet text."""
    add_shape(slide, left, top, width, height, fill_color=CARD_BG, line_color=accent, line_width=Pt(1.5))
    # icon circle
    if icon_text:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3), Inches(0.7), Inches(0.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = accent
        circ.line.fill.background()
        circ.text_frame.paragraphs[0].text = icon_text
        circ.text_frame.paragraphs[0].font.size = Pt(22)
        circ.text_frame.paragraphs[0].font.color.rgb = WHITE
        circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circ.text_frame.paragraphs[0].font.name = "Segoe UI Emoji"

    title_left = left + (Inches(1.2) if icon_text else Inches(0.3))
    t = tb(slide, title_left, top + Inches(0.3), width - Inches(1.5), Inches(0.5))
    set_text(t, title, size=20, color=WHITE, bold=True)

    body_top = top + Inches(1.1)
    for i, line in enumerate(body_lines):
        t = tb(slide, left + Inches(0.4), body_top + Inches(i * 0.38), width - Inches(0.8), Inches(0.4))
        set_text(t, f"•  {line}", size=14, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
# Top accent line
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_BLUE)

t = tb(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5))
set_text(t, "IoT-Powered Smart Manufacturing", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

t = tb(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.2))
set_text(t, "OEE Measurement  •  Supply Chain Automation  •  Real-Time Goods Tracking", size=26, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

t = tb(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8))
set_text(t, "Transforming Production Visibility with Real-Time Data Intelligence", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, Inches(4), Inches(6.0), Inches(5.333), Pt(2), ACCENT_BLUE)

t = tb(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5))
set_text(t, "Confidential  |  Prepared for Client Review  |  March 2026", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "Agenda", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.15), Inches(2), Pt(4), ACCENT_BLUE)

agenda_items = [
    ("01", "The Challenge — Why OEE & Supply Chain Visibility Matter", ACCENT_BLUE),
    ("02", "Our Solution — IoT-Powered Real-Time OEE Platform", ACCENT_GREEN),
    ("03", "Live Dashboard Demo — Production Machine Analytics", ACCENT_ORANGE),
    ("04", "OEE Deep Dive — Availability, Performance, Quality", ACCENT_BLUE),
    ("05", "Supply Chain Management & Communication", ACCENT_GREEN),
    ("06", "Automated Tracking of Goods — End to End", ACCENT_ORANGE),
    ("07", "System Architecture & Technology Stack", ACCENT_BLUE),
    ("08", "ROI & Business Impact", ACCENT_GREEN),
    ("09", "Implementation Roadmap", ACCENT_ORANGE),
    ("10", "Next Steps", ACCENT_BLUE),
]

for i, (num, text, accent) in enumerate(agenda_items):
    y = Inches(1.6) + Inches(i * 0.55)
    # number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y, Inches(0.7), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    badge.text_frame.paragraphs[0].text = num
    badge.text_frame.paragraphs[0].font.size = Pt(16)
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.paragraphs[0].font.name = "Calibri"
    # text
    t = tb(slide, Inches(2.0), y, Inches(10), Inches(0.45))
    set_text(t, text, size=20, color=WHITE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE CHALLENGE
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "The Challenge", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.15), Inches(2.5), Pt(4), ACCENT_ORANGE)

# Left column — Manufacturing
card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.6),
     "Manufacturing Floor", [
         "No real-time visibility into machine performance",
         "Downtime detected too late — hours of lost production",
         "Manual data collection → errors & delays",
         "OEE calculated weekly/monthly, not in real-time",
     ], "🏭", ACCENT_ORANGE)

# Right column — Supply Chain
card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.6),
     "Supply Chain & Logistics", [
         "Goods tracking relies on manual check-ins",
         "No automated communication between nodes",
         "Lost/delayed shipments discovered too late",
         "Siloed data across warehouse, transport, delivery",
     ], "📦", ACCENT_BLUE)

# Bottom impact bar
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), fill_color=CARD_BG, line_color=RGBColor(0xFF, 0x45, 0x45), line_width=Pt(2))
t = tb(slide, Inches(0.8), Inches(4.9), Inches(12), Inches(0.5))
set_text(t, "💰  Business Impact of Inaction", size=24, color=RGBColor(0xFF, 0x45, 0x45), bold=True)

impacts = [
    "Average manufacturer operates at only 60% OEE — world-class is 85%+",
    "Unplanned downtime costs $260K/hour on average (Aberdeen Group)",
    "30% of goods in transit lack real-time visibility (McKinsey)",
    "Supply chain disruptions cost companies 45% of one year's profits over a decade (McKinsey)",
]
for i, imp in enumerate(impacts):
    t = tb(slide, Inches(1.2), Inches(5.45) + Inches(i * 0.38), Inches(11), Inches(0.4))
    set_text(t, f"▸  {imp}", size=15, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR SOLUTION OVERVIEW
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "Our Solution — IoT-Powered Intelligence Platform", size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4), Pt(4), ACCENT_GREEN)

# Three pillars
pillars = [
    ("📊", "Real-Time OEE", "Measure Availability, Performance\n& Quality — live, every second", ACCENT_BLUE),
    ("🔗", "Supply Chain Hub", "Automated communication between\nall supply chain nodes", ACCENT_GREEN),
    ("📍", "Goods Tracking", "End-to-end automated tracking\nfrom production to delivery", ACCENT_ORANGE),
]

for i, (icon, title, desc, accent) in enumerate(pillars):
    x = Inches(0.5) + Inches(i * 4.2)
    add_shape(slide, x, Inches(1.6), Inches(3.8), Inches(2.8), fill_color=CARD_BG, line_color=accent, line_width=Pt(2))
    # Icon
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.4), Inches(1.9), Inches(1), Inches(1))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].text = icon
    circ.text_frame.paragraphs[0].font.size = Pt(32)
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circ.text_frame.paragraphs[0].font.name = "Segoe UI Emoji"
    # Title
    t = tb(slide, x + Inches(0.2), Inches(3.1), Inches(3.4), Inches(0.5))
    set_text(t, title, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Desc
    t = tb(slide, x + Inches(0.2), Inches(3.6), Inches(3.4), Inches(0.7))
    set_text(t, desc, size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Key value prop
add_shape(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.2), fill_color=CARD_BG, line_color=HIGHLIGHT, line_width=Pt(2))
t = tb(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(0.5))
set_text(t, "✨  Unified Platform — One Dashboard for Everything", size=24, color=HIGHLIGHT, bold=True)

values = [
    "Sensors on machines → real-time pulse data → automatic OEE calculation",
    "IoT gateways at each supply chain node → automatic status updates",
    "RFID/BLE/GPS tags on goods → live location tracking on map",
    "Automated alerts via SMS, Email, WhatsApp when KPIs breach thresholds",
]
for i, v in enumerate(values):
    t = tb(slide, Inches(1.2), Inches(5.5) + Inches(i * 0.38), Inches(11), Inches(0.4))
    set_text(t, f"✓  {v}", size=15, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LIVE DASHBOARD DEMO
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "Live Dashboard — Production Machine Analytics", size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4), Pt(4), ACCENT_BLUE)

# Dashboard mockup panels
panels = [
    (Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), "Live Pulse Stream", "Real-time sensor data from production machine\n\nEvery pulse captured as it happens\nGaps between clusters = material feed time"),
    (Inches(6.8), Inches(1.5), Inches(3), Inches(1.1), "Avg Pulse Gap\n52.8 ms", ""),
    (Inches(10.1), Inches(1.5), Inches(2.7), Inches(1.1), "Avg Feed Gap\n505 ms", ""),
    (Inches(6.8), Inches(2.8), Inches(3), Inches(1.1), "Products/5min\n47", ""),
    (Inches(10.1), Inches(2.8), Inches(2.7), Inches(1.1), "Total Pulses\n534", ""),
    (Inches(0.5), Inches(4.3), Inches(6), Inches(2.5), "Key Metrics Over Time", "Trending graph showing:\n• Avg Intra-Series Gap\n• Avg Inter-Series Gap\n• Products Produced\n• Total Pulses"),
    (Inches(6.8), Inches(4.3), Inches(6), Inches(2.5), "Gap Distribution & Histogram", "• Pulses per product (bar chart)\n• All consecutive gaps (line chart)\n• Low values = production pulses\n• High spikes = material feed"),
]

for (x, y, w, h, title, desc) in panels:
    add_shape(slide, x, y, w, h, fill_color=CARD_BG, line_color=ACCENT_BLUE, line_width=Pt(1))
    t = tb(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.4))
    set_text(t, title, size=16, color=ACCENT_BLUE, bold=True)
    if desc:
        t = tb(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), h - Inches(0.6))
        set_text(t, desc, size=13, color=LIGHT_GRAY)

# Live URL note
t = tb(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4))
set_text(t, "🔴 LIVE: https://fe055024.grafana.net  —  Auto-refreshes every 5 seconds", size=16, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — OEE DEEP DIVE
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "OEE — Overall Equipment Effectiveness", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.15), Inches(4), Pt(4), ACCENT_BLUE)

# OEE Formula
add_shape(slide, Inches(2), Inches(1.6), Inches(9.3), Inches(1.0), fill_color=CARD_BG, line_color=ACCENT_BLUE, line_width=Pt(2))
t = tb(slide, Inches(2.2), Inches(1.7), Inches(9), Inches(0.7))
set_text(t, "OEE  =  Availability  ×  Performance  ×  Quality", size=32, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

# Three components
oee_items = [
    ("Availability", "Is the machine running?", [
        "Planned production time vs actual run time",
        "Detects unplanned stops automatically",
        "Measured via inter-series gaps (>200ms)",
        "Target: ≥ 90%"
    ], ACCENT_GREEN, "⏱"),
    ("Performance", "Is it running at full speed?", [
        "Ideal cycle time vs actual cycle time",
        "Detects slow cycles via intra-series gaps",
        "Measured via pulse timing (avg 52ms ideal)",
        "Target: ≥ 95%"
    ], ACCENT_BLUE, "⚡"),
    ("Quality", "Are products defect-free?", [
        "Good products vs total products",
        "Detects anomalous pulse patterns",
        "Irregular pulse count = potential defect",
        "Target: ≥ 99%"
    ], ACCENT_ORANGE, "✅"),
]

for i, (title, subtitle, items, accent, icon) in enumerate(oee_items):
    x = Inches(0.5) + Inches(i * 4.2)
    add_shape(slide, x, Inches(3.0), Inches(3.9), Inches(4.1), fill_color=CARD_BG, line_color=accent, line_width=Pt(2))
    # Icon
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), Inches(3.2), Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].text = icon
    circ.text_frame.paragraphs[0].font.size = Pt(20)
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title & subtitle
    t = tb(slide, x + Inches(0.9), Inches(3.15), Inches(2.8), Inches(0.4))
    set_text(t, title, size=24, color=WHITE, bold=True)
    t = tb(slide, x + Inches(0.9), Inches(3.55), Inches(2.8), Inches(0.35))
    set_text(t, subtitle, size=14, color=accent)
    # Bullets
    for j, item in enumerate(items):
        t = tb(slide, x + Inches(0.25), Inches(4.1) + Inches(j * 0.4), Inches(3.5), Inches(0.4))
        set_text(t, f"•  {item}", size=13, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — HOW WE MEASURE OEE (IoT Detail)
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "How We Measure OEE — From Sensor to Dashboard", size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4.5), Pt(4), ACCENT_GREEN)

# Flow diagram as boxes with arrows
steps = [
    ("🔌", "Sensor", "Proximity/pulse sensor\non machine spindle", ACCENT_BLUE),
    ("→", "", "", WHITE),
    ("📡", "IoT Gateway", "Edge device captures\npulse timestamps", ACCENT_GREEN),
    ("→", "", "", WHITE),
    ("🗄", "InfluxDB", "Time-series database\nstores every pulse", ACCENT_ORANGE),
    ("→", "", "", WHITE),
    ("📊", "Grafana", "Real-time dashboard\ncalculates OEE live", ACCENT_BLUE),
    ("→", "", "", WHITE),
    ("🔔", "Alerts", "SMS/Email/WhatsApp\nwhen KPI breaches", RGBColor(0xFF, 0x45, 0x45)),
]

x = Inches(0.3)
for icon, title, desc, accent in steps:
    if icon == "→":
        t = tb(slide, x, Inches(2.3), Inches(0.5), Inches(0.5))
        set_text(t, "→", size=30, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        x += Inches(0.5)
    else:
        w = Inches(1.5)
        add_shape(slide, x, Inches(1.7), w, Inches(2.0), fill_color=CARD_BG, line_color=accent, line_width=Pt(1.5))
        t = tb(slide, x, Inches(1.8), w, Inches(0.5))
        set_text(t, icon, size=28, align=PP_ALIGN.CENTER)
        t = tb(slide, x + Inches(0.05), Inches(2.3), w - Inches(0.1), Inches(0.4))
        set_text(t, title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        t = tb(slide, x + Inches(0.05), Inches(2.7), w - Inches(0.1), Inches(0.8))
        set_text(t, desc, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        x += Inches(1.6)

# Key metrics table
add_shape(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.0), fill_color=CARD_BG, line_color=ACCENT_BLUE, line_width=Pt(1))
t = tb(slide, Inches(0.8), Inches(4.2), Inches(12), Inches(0.5))
set_text(t, "What Our Sensors Capture in Real-Time", size=22, color=ACCENT_BLUE, bold=True)

metrics = [
    ("Metric", "How It's Measured", "OEE Component", "Dashboard Panel"),
    ("Intra-Series Gap", "Time between consecutive pulses (<200ms)", "Performance", "Avg Pulse Gap stat + trend"),
    ("Inter-Series Gap", "Time between products (≥200ms = material feed)", "Availability", "Avg Feed Gap stat + trend"),
    ("Pulse Count/Product", "Number of pulses in each series", "Quality", "Pulses Per Product histogram"),
    ("Products/Minute", "Count of inter-series boundaries", "Availability", "Products Produced stat + trend"),
    ("Total Throughput", "All pulses counted over time", "Performance", "Total Pulses + Live Stream"),
]

for i, row in enumerate(metrics):
    y = Inches(4.75) + Inches(i * 0.36)
    color = ACCENT_BLUE if i == 0 else LIGHT_GRAY
    bold = i == 0
    for j, cell in enumerate(row):
        x = Inches(0.8) + Inches(j * 2.9)
        t = tb(slide, x, y, Inches(2.8), Inches(0.35))
        set_text(t, cell, size=12 if i > 0 else 13, color=color, bold=bold)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION DIVIDER — SUPPLY CHAIN
# ═════════════════════════════════════════════════════════════════════════════
section_divider(
    "Supply Chain Management",
    "Communication  •  Automation  •  Real-Time Goods Tracking",
    ACCENT_GREEN
)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SUPPLY CHAIN OVERVIEW
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "Supply Chain Communication & Automation", size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4.5), Pt(4), ACCENT_GREEN)

# Supply chain flow
sc_nodes = [
    ("�icing", "Raw Material\nSupplier", ACCENT_ORANGE),
    ("🏭", "Manufacturing\nPlant", ACCENT_BLUE),
    ("🏪", "Warehouse\n/ DC", ACCENT_GREEN),
    ("🚛", "Transport\n& Logistics", ACCENT_ORANGE),
    ("🏬", "Retail /\nEnd Customer", ACCENT_BLUE),
]
# Fix first icon
sc_nodes[0] = ("📋", "Raw Material\nSupplier", ACCENT_ORANGE)

for i, (icon, label, accent) in enumerate(sc_nodes):
    x = Inches(0.4) + Inches(i * 2.6)
    add_shape(slide, x, Inches(1.6), Inches(2.2), Inches(1.5), fill_color=CARD_BG, line_color=accent, line_width=Pt(2))
    t = tb(slide, x, Inches(1.7), Inches(2.2), Inches(0.5))
    set_text(t, icon, size=30, align=PP_ALIGN.CENTER)
    t = tb(slide, x + Inches(0.1), Inches(2.2), Inches(2), Inches(0.6))
    set_text(t, label, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Arrow
    if i < len(sc_nodes) - 1:
        t = tb(slide, x + Inches(2.15), Inches(2.0), Inches(0.5), Inches(0.5))
        set_text(t, "→", size=28, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# Communication features below
t = tb(slide, Inches(0.8), Inches(3.4), Inches(12), Inches(0.5))
set_text(t, "Automated Communication at Every Node", size=24, color=ACCENT_GREEN, bold=True)

comm_features = [
    ("📱 Real-Time Alerts", "Automated SMS, Email & WhatsApp notifications\nwhen goods arrive, depart, or are delayed at any node", Inches(0.5), ACCENT_BLUE),
    ("🔄 EDI / API Integration", "Automated purchase orders, invoices, and ASNs\nexchanged between supplier ↔ manufacturer ↔ warehouse", Inches(4.5), ACCENT_GREEN),
    ("📊 Unified Dashboard", "Single pane of glass for all supply chain KPIs:\nInventory levels, transit times, delivery accuracy", Inches(8.5), ACCENT_ORANGE),
]

for title, desc, x, accent in comm_features:
    add_shape(slide, x, Inches(3.95), Inches(3.8), Inches(2.2), fill_color=CARD_BG, line_color=accent, line_width=Pt(1.5))
    t = tb(slide, x + Inches(0.2), Inches(4.05), Inches(3.4), Inches(0.4))
    set_text(t, title, size=16, color=WHITE, bold=True)
    t = tb(slide, x + Inches(0.2), Inches(4.5), Inches(3.4), Inches(1.3))
    set_text(t, desc, size=13, color=LIGHT_GRAY)

# Bottom key stat
add_shape(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.8), fill_color=CARD_BG, line_color=HIGHLIGHT, line_width=Pt(2))
t = tb(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.5))
set_text(t, "🎯 Result: 90% reduction in manual communication  |  75% faster issue resolution  |  99.5% notification delivery", size=16, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — AUTOMATED GOODS TRACKING
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "Automated Tracking of Goods — End to End", size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4), Pt(4), ACCENT_ORANGE)

# Tracking technologies
tech_cards = [
    ("📻", "RFID Tags", "Passive/active RFID on every\nproduct, pallet & container.\nAuto-scanned at gates.", ACCENT_BLUE),
    ("📶", "BLE Beacons", "Bluetooth Low Energy for\nindoor warehouse tracking.\nSub-meter accuracy.", ACCENT_GREEN),
    ("🛰", "GPS Tracking", "Real-time vehicle & shipment\nlocation during transit.\nGeofence alerts.", ACCENT_ORANGE),
    ("📷", "Vision / Barcode", "Camera-based scanning at\nloading/unloading docks.\nAI defect detection.", RGBColor(0xBB, 0x66, 0xFF)),
]

for i, (icon, title, desc, accent) in enumerate(tech_cards):
    x = Inches(0.3) + Inches(i * 3.3)
    add_shape(slide, x, Inches(1.5), Inches(3.0), Inches(2.5), fill_color=CARD_BG, line_color=accent, line_width=Pt(2))
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.0), Inches(1.7), Inches(0.8), Inches(0.8))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].text = icon
    circ.text_frame.paragraphs[0].font.size = Pt(26)
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    t = tb(slide, x + Inches(0.15), Inches(2.6), Inches(2.7), Inches(0.4))
    set_text(t, title, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    t = tb(slide, x + Inches(0.15), Inches(3.05), Inches(2.7), Inches(0.9))
    set_text(t, desc, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Tracking journey
add_shape(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), fill_color=CARD_BG, line_color=ACCENT_ORANGE, line_width=Pt(1.5))
t = tb(slide, Inches(0.8), Inches(4.4), Inches(12), Inches(0.5))
set_text(t, "🔍  Automated Tracking Journey — No Human Intervention Required", size=20, color=ACCENT_ORANGE, bold=True)

journey_steps = [
    ("1", "PRODUCED", "Product tagged with\nRFID at production line", ACCENT_BLUE),
    ("2", "PACKED", "Auto-scanned at\npacking station", ACCENT_GREEN),
    ("3", "STORED", "BLE tracks location\nin warehouse", ACCENT_ORANGE),
    ("4", "SHIPPED", "GPS tracks vehicle\nin real-time", ACCENT_BLUE),
    ("5", "DELIVERED", "Auto-confirmed at\ncustomer dock", ACCENT_GREEN),
]

for i, (num, title, desc, accent) in enumerate(journey_steps):
    x = Inches(0.7) + Inches(i * 2.45)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), Inches(5.0), Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].text = num
    circ.text_frame.paragraphs[0].font.size = Pt(20)
    circ.text_frame.paragraphs[0].font.color.rgb = WHITE
    circ.text_frame.paragraphs[0].font.bold = True
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title
    t = tb(slide, x, Inches(5.7), Inches(1.8), Inches(0.35))
    set_text(t, title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Desc
    t = tb(slide, x, Inches(6.05), Inches(1.8), Inches(0.7))
    set_text(t, desc, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # Arrow
    if i < 4:
        t = tb(slide, x + Inches(1.85), Inches(5.1), Inches(0.6), Inches(0.4))
        set_text(t, "→", size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ARCHITECTURE
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "System Architecture & Technology Stack", size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4), Pt(4), ACCENT_BLUE)

# Architecture layers
layers = [
    ("EDGE LAYER", "Sensors, RFID, BLE, GPS\nIoT Gateways, Edge Computing", ACCENT_GREEN, Inches(1.5)),
    ("INGESTION LAYER", "MQTT / HTTP → InfluxDB\nReal-time time-series storage", ACCENT_BLUE, Inches(2.8)),
    ("ANALYTICS LAYER", "Flux Queries, OEE Calculation\nAnomaly Detection, ML Pipeline", ACCENT_ORANGE, Inches(4.1)),
    ("VISUALIZATION LAYER", "Grafana Cloud Dashboard\nMobile App, API Endpoints", RGBColor(0xBB, 0x66, 0xFF), Inches(5.4)),
]

for title, desc, accent, y in layers:
    add_shape(slide, Inches(0.5), y, Inches(7.5), Inches(1.1), fill_color=CARD_BG, line_color=accent, line_width=Pt(2))
    t = tb(slide, Inches(0.8), y + Inches(0.1), Inches(3), Inches(0.4))
    set_text(t, title, size=16, color=accent, bold=True)
    t = tb(slide, Inches(0.8), y + Inches(0.45), Inches(7), Inches(0.5))
    set_text(t, desc, size=13, color=LIGHT_GRAY)

# Tech stack sidebar
add_shape(slide, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.0), fill_color=CARD_BG, line_color=ACCENT_BLUE, line_width=Pt(1))
t = tb(slide, Inches(8.5), Inches(1.6), Inches(4), Inches(0.5))
set_text(t, "Technology Stack", size=20, color=ACCENT_BLUE, bold=True)

stack = [
    "Database: InfluxDB 2.7 (Time-Series)",
    "Dashboard: Grafana Cloud (10.x)",
    "Query: Flux (real-time analytics)",
    "IoT Protocol: MQTT, HTTP/REST",
    "Edge: Python, Node.js, C/Arduino",
    "Tracking: RFID, BLE, GPS",
    "Alerts: Grafana Alerting + Webhooks",
    "Hosting: Docker, Cloud-native",
    "Security: TLS, Token-based auth",
    "Integration: REST API, EDI, Webhooks",
]
for i, item in enumerate(stack):
    t = tb(slide, Inches(8.6), Inches(2.2) + Inches(i * 0.4), Inches(4), Inches(0.35))
    set_text(t, f"▸  {item}", size=13, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ROI & BUSINESS IMPACT
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "ROI & Business Impact", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.15), Inches(3), Pt(4), ACCENT_GREEN)

# KPI cards
kpis = [
    ("25%", "OEE\nImprovement", "From 60% to 85%+\nwithin 6 months", ACCENT_GREEN),
    ("40%", "Downtime\nReduction", "Predictive alerts\ncatch issues early", ACCENT_BLUE),
    ("90%", "Manual Work\nEliminated", "Automated tracking\n& communication", ACCENT_ORANGE),
    ("3x", "Faster\nDecisions", "Real-time data vs\nweekly reports", RGBColor(0xBB, 0x66, 0xFF)),
]

for i, (number, title, desc, accent) in enumerate(kpis):
    x = Inches(0.3) + Inches(i * 3.3)
    add_shape(slide, x, Inches(1.6), Inches(3.0), Inches(2.5), fill_color=CARD_BG, line_color=accent, line_width=Pt(2))
    t = tb(slide, x, Inches(1.8), Inches(3.0), Inches(0.8))
    set_text(t, number, size=48, color=accent, bold=True, align=PP_ALIGN.CENTER)
    t = tb(slide, x + Inches(0.2), Inches(2.5), Inches(2.6), Inches(0.5))
    set_text(t, title, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    t = tb(slide, x + Inches(0.2), Inches(3.0), Inches(2.6), Inches(0.6))
    set_text(t, desc, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ROI timeline
add_shape(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7), fill_color=CARD_BG, line_color=ACCENT_GREEN, line_width=Pt(1.5))
t = tb(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(0.5))
set_text(t, "💰  Estimated ROI Timeline", size=22, color=ACCENT_GREEN, bold=True)

roi_rows = [
    ("Timeline", "Investment", "Savings", "Net Impact"),
    ("Month 1-3", "Setup & deployment", "10-15% downtime reduction", "Foundation laid"),
    ("Month 4-6", "Tuning & optimization", "20-25% OEE improvement", "Breaking even"),
    ("Month 7-12", "Scale & automate", "40% cost reduction in tracking", "2-3x ROI"),
    ("Year 2+", "Predictive analytics", "Continuous improvement", "5-10x ROI"),
]
for i, row in enumerate(roi_rows):
    y = Inches(5.1) + Inches(i * 0.38)
    bold = i == 0
    color = ACCENT_BLUE if i == 0 else LIGHT_GRAY
    for j, cell in enumerate(row):
        x = Inches(1.0) + Inches(j * 2.8)
        t = tb(slide, x, y, Inches(2.7), Inches(0.35))
        set_text(t, cell, size=13, color=color, bold=bold)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — IMPLEMENTATION ROADMAP
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t, "Implementation Roadmap", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.15), Inches(3.5), Pt(4), ACCENT_ORANGE)

phases = [
    ("Phase 1\nWeeks 1-4", "Foundation", [
        "Install sensors on pilot machines",
        "Deploy InfluxDB + Grafana Cloud",
        "Configure basic OEE dashboard",
        "Set up alerting rules",
    ], ACCENT_BLUE),
    ("Phase 2\nWeeks 5-8", "Supply Chain Integration", [
        "Deploy RFID/BLE at key nodes",
        "Integrate supplier communication",
        "Automated PO/ASN exchange",
        "Goods tracking dashboard",
    ], ACCENT_GREEN),
    ("Phase 3\nWeeks 9-12", "Automation & Scale", [
        "Roll out to all machines",
        "GPS tracking for logistics",
        "Automated reporting & escalation",
        "Mobile app for field teams",
    ], ACCENT_ORANGE),
    ("Phase 4\nOngoing", "Intelligence", [
        "ML-based anomaly detection",
        "Predictive maintenance alerts",
        "Demand forecasting integration",
        "Continuous improvement loop",
    ], RGBColor(0xBB, 0x66, 0xFF)),
]

for i, (phase, title, items, accent) in enumerate(phases):
    x = Inches(0.3) + Inches(i * 3.3)
    add_shape(slide, x, Inches(1.6), Inches(3.0), Inches(5.3), fill_color=CARD_BG, line_color=accent, line_width=Pt(2))
    # Phase label
    t = tb(slide, x + Inches(0.2), Inches(1.7), Inches(2.6), Inches(0.8))
    set_text(t, phase, size=14, color=accent, bold=True, align=PP_ALIGN.CENTER)
    # Title
    t = tb(slide, x + Inches(0.2), Inches(2.5), Inches(2.6), Inches(0.4))
    set_text(t, title, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Divider
    add_rect(slide, x + Inches(0.8), Inches(2.95), Inches(1.4), Pt(2), accent)
    # Items
    for j, item in enumerate(items):
        t = tb(slide, x + Inches(0.2), Inches(3.2) + Inches(j * 0.55), Inches(2.6), Inches(0.5))
        set_text(t, f"✓  {item}", size=13, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — NEXT STEPS & CTA
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

t = tb(slide, Inches(1), Inches(0.8), Inches(11), Inches(1))
set_text(t, "Next Steps", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(1.7), Inches(2.3), Pt(4), ACCENT_BLUE)

steps = [
    ("1", "Pilot Program", "Start with 1-2 machines for OEE monitoring\nand 1 supply chain route for tracking", ACCENT_BLUE),
    ("2", "Proof of Value", "4-week pilot to demonstrate measurable\nOEE improvement & tracking accuracy", ACCENT_GREEN),
    ("3", "Scale & Deploy", "Roll out across all production lines\nand supply chain nodes", ACCENT_ORANGE),
]

for i, (num, title, desc, accent) in enumerate(steps):
    y = Inches(2.2) + Inches(i * 1.6)
    # Number
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), y + Inches(0.1), Inches(0.8), Inches(0.8))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].text = num
    circ.text_frame.paragraphs[0].font.size = Pt(28)
    circ.text_frame.paragraphs[0].font.color.rgb = WHITE
    circ.text_frame.paragraphs[0].font.bold = True
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Content
    t = tb(slide, Inches(3.2), y, Inches(7.5), Inches(0.5))
    set_text(t, title, size=26, color=WHITE, bold=True)
    t = tb(slide, Inches(3.2), y + Inches(0.5), Inches(7.5), Inches(0.8))
    set_text(t, desc, size=16, color=LIGHT_GRAY)

# Bottom CTA
add_shape(slide, Inches(3), Inches(6.2), Inches(7.3), Inches(0.8), fill_color=ACCENT_BLUE, line_color=ACCENT_BLUE, line_width=Pt(0))
t = tb(slide, Inches(3), Inches(6.3), Inches(7.3), Inches(0.6))
set_text(t, "📞  Let's Schedule a Live Demo & Pilot Discussion", size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

t = tb(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5))
set_text(t, "Thank You", size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

t = tb(slide, Inches(1), Inches(3.5), Inches(11), Inches(1))
set_text(t, "Real-Time OEE  •  Smart Supply Chain  •  Automated Tracking", size=24, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(5), Inches(4.5), Inches(3.3), Pt(2), ACCENT_BLUE)

t = tb(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6))
set_text(t, "Live Dashboard: https://fe055024.grafana.net", size=18, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

t = tb(slide, Inches(1), Inches(5.7), Inches(11), Inches(0.6))
set_text(t, "Questions? Let's discuss.", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SAVE
# ═════════════════════════════════════════════════════════════════════════════
output_path = "/home/ashok/IOT/IoT_OEE_SupplyChain_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
